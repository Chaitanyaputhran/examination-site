#!/usr/bin/env python3
"""
Generate a concise, beautiful PowerPoint presentation for Examination Portal
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def add_styled_title_slide(prs, title, subtitle):
    """Add a beautifully styled title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background gradient effect
    bg1 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(7.5)
    )
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = RGBColor(15, 76, 129)
    bg1.line.fill.background()
    slide.shapes._spTree.remove(bg1._element)
    slide.shapes._spTree.insert(2, bg1._element)

    # Decorative circles
    circle1 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(-2), Inches(-2),
        Inches(6), Inches(6)
    )
    circle1.fill.solid()
    circle1.fill.fore_color.rgb = RGBColor(26, 115, 186)
    circle1.line.fill.background()

    circle2 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(7), Inches(4),
        Inches(5), Inches(5)
    )
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = RGBColor(59, 130, 246)
    circle2.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(72)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = RGBColor(200, 230, 255)
    subtitle_para.alignment = PP_ALIGN.CENTER

def add_tech_stack_slide(prs):
    """Add tech stack slide with two columns"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header bar
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(1.4)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(220, 38, 38)
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "🛠️ Technology Stack"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)

    # Backend column background
    backend_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.8),
        Inches(4.5), Inches(5.2)
    )
    backend_bg.fill.solid()
    backend_bg.fill.fore_color.rgb = RGBColor(254, 242, 242)
    backend_bg.line.color.rgb = RGBColor(220, 38, 38)
    backend_bg.line.width = Pt(3)

    # Frontend column background
    frontend_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.2), Inches(1.8),
        Inches(4.5), Inches(5.2)
    )
    frontend_bg.fill.solid()
    frontend_bg.fill.fore_color.rgb = RGBColor(239, 246, 255)
    frontend_bg.line.color.rgb = RGBColor(59, 130, 246)
    frontend_bg.line.width = Pt(3)

    # Backend title
    backend_title = slide.shapes.add_textbox(Inches(0.7), Inches(2), Inches(4), Inches(0.5))
    bt_frame = backend_title.text_frame
    bt_frame.text = "⚙️ Backend"
    bt_para = bt_frame.paragraphs[0]
    bt_para.font.size = Pt(32)
    bt_para.font.bold = True
    bt_para.font.color.rgb = RGBColor(220, 38, 38)
    bt_para.alignment = PP_ALIGN.CENTER

    # Frontend title
    frontend_title = slide.shapes.add_textbox(Inches(5.4), Inches(2), Inches(4), Inches(0.5))
    ft_frame = frontend_title.text_frame
    ft_frame.text = "⚛️ Frontend"
    ft_para = ft_frame.paragraphs[0]
    ft_para.font.size = Pt(32)
    ft_para.font.bold = True
    ft_para.font.color.rgb = RGBColor(59, 130, 246)
    ft_para.alignment = PP_ALIGN.CENTER

    # Backend content
    backend_content = [
        "☕ Java 17",
        "🍃 Spring Boot 3.2.0",
        "  • Spring Security + JWT",
        "  • Spring Data JPA",
        "  • Spring Mail",
        "",
        "🗄️ MySQL 8.0",
        "📦 Maven",
        "📄 iText7 (PDF)",
        "🔒 BCrypt Encryption"
    ]

    backend_box = slide.shapes.add_textbox(Inches(0.9), Inches(2.7), Inches(4), Inches(4.1))
    backend_frame = backend_box.text_frame
    backend_frame.word_wrap = True
    for item in backend_content:
        p = backend_frame.add_paragraph()
        p.text = item
        if item.startswith("  •"):
            p.font.size = Pt(16)
            p.level = 1
        else:
            p.font.size = Pt(20)
            p.font.bold = not item.startswith("  ")
        p.font.color.rgb = RGBColor(30, 41, 59)
        p.space_before = Pt(8)

    # Frontend content
    frontend_content = [
        "⚛️ React 18.2.0",
        "🎨 Tailwind CSS 3.3.6",
        "🔀 React Router 6.20.0",
        "📡 Axios",
        "📊 Recharts 2.10.3",
        "🔔 React Hot Toast",
        "⚡ Vite 5.0.8",
        "📦 npm"
    ]

    frontend_box = slide.shapes.add_textbox(Inches(5.4), Inches(2.7), Inches(4), Inches(4.1))
    frontend_frame = frontend_box.text_frame
    frontend_frame.word_wrap = True
    for item in frontend_content:
        p = frontend_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(30, 41, 59)
        p.space_before = Pt(10)

def add_features_slide(prs, slide_num):
    """Add features slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    if slide_num == 1:
        header_color = (16, 185, 129)
        header_text = "✨ Key Features"
        features = [
            ("🔐 Secure Authentication", "JWT-based authentication with role-based access control (Admin & Student)"),
            ("📝 Test Management", "Create tests with MCQ questions, set duration, marks, and passing criteria"),
            ("⚡ Real-time Exams", "Live countdown timer, auto-save answers, and automatic submission"),
            ("🔄 Unlimited Retakes", "Students can retake any test unlimited times to improve their scores"),
            ("📊 Rich Analytics", "Comprehensive dashboard with user, test, and performance statistics"),
            ("📧 Email Reports", "Automated PDF generation and email delivery of detailed test results")
        ]
    else:
        header_color = (139, 92, 246)
        header_text = "🎯 Advanced Capabilities"
        features = [
            ("👥 User Management", "Create and manage admin and student accounts with activation controls"),
            ("📚 Subject Organization", "Organize tests by subjects/categories for better content management"),
            ("📈 Performance Tracking", "Subject-wise analytics, attempt history, and score comparisons"),
            ("✍️ Flexible Navigation", "Navigate freely between questions, flag for review during exams"),
            ("🎨 Responsive Design", "Modern, mobile-friendly interface built with Tailwind CSS"),
            ("🛡️ Data Security", "BCrypt password hashing, SQL injection prevention, CORS protection")
        ]

    # Header bar
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(1.4)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(*header_color)
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = header_text
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)

    # Feature cards
    positions = [
        (0.5, 1.8), (5.2, 1.8),
        (0.5, 4.1), (5.2, 4.1),
        (0.5, 6.4), (5.2, 6.4)
    ]

    colors = [
        (239, 246, 255, 59, 130, 246),   # Blue
        (236, 254, 255, 6, 182, 212),    # Cyan
        (240, 253, 244, 34, 197, 94),    # Green
        (254, 249, 195, 234, 179, 8),    # Yellow
        (254, 242, 242, 239, 68, 68),    # Red
        (243, 232, 255, 168, 85, 247)    # Purple
    ]

    for i, ((x, y), (feat_title, feat_desc)) in enumerate(zip(positions, features)):
        bg_color, border_color = colors[i][:3], colors[i][3:]

        # Card background
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y),
            Inches(4.5), Inches(2.1)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(*bg_color)
        card.line.color.rgb = RGBColor(*border_color)
        card.line.width = Pt(2.5)

        # Feature title
        title_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.15), Inches(4.1), Inches(0.5))
        title_frame = title_box.text_frame
        title_frame.text = feat_title
        title_frame.word_wrap = True
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(22)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(*border_color)

        # Feature description
        desc_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.7), Inches(4.1), Inches(1.2))
        desc_frame = desc_box.text_frame
        desc_frame.word_wrap = True
        desc_frame.text = feat_desc
        desc_para = desc_frame.paragraphs[0]
        desc_para.font.size = Pt(15)
        desc_para.font.color.rgb = RGBColor(51, 65, 85)
        desc_para.line_spacing = 1.2

def add_thank_you_slide(prs):
    """Add thank you slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(15, 76, 129)
    bg.line.fill.background()
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    # Decorative circles
    for x, y, size in [(-1, -1, 4), (8, 5, 5)]:
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x), Inches(y),
            Inches(size), Inches(size)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = RGBColor(26, 115, 186)
        circle.line.fill.background()

    # Thank you text
    thank_you_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(6), Inches(1))
    thank_you_frame = thank_you_box.text_frame
    thank_you_frame.text = "Thank You!"
    thank_you_para = thank_you_frame.paragraphs[0]
    thank_you_para.font.size = Pt(72)
    thank_you_para.font.bold = True
    thank_you_para.font.color.rgb = RGBColor(255, 255, 255)
    thank_you_para.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(2), Inches(4.2), Inches(6), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Examination Portal"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(32)
    subtitle_para.font.color.rgb = RGBColor(200, 230, 255)
    subtitle_para.alignment = PP_ALIGN.CENTER

    # Footer
    footer_box = slide.shapes.add_textbox(Inches(2), Inches(5.3), Inches(6), Inches(0.5))
    footer_frame = footer_box.text_frame
    footer_frame.text = "Modern Online Examination Management System"
    footer_para = footer_frame.paragraphs[0]
    footer_para.font.size = Pt(18)
    footer_para.font.color.rgb = RGBColor(200, 230, 255)
    footer_para.alignment = PP_ALIGN.CENTER

def create_concise_presentation():
    """Create a concise, beautiful presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_styled_title_slide(prs,
                          "Examination Portal",
                          "Modern Online Examination Management System")

    # Slide 2: Tech Stack
    add_tech_stack_slide(prs)

    # Slide 3: Features Part 1
    add_features_slide(prs, 1)

    # Slide 4: Features Part 2
    add_features_slide(prs, 2)

    # Slide 5: Thank You
    add_thank_you_slide(prs)

    # Save presentation
    output_file = "Examination_Portal.pptx"
    prs.save(output_file)
    print(f"✅ Concise presentation created: {output_file}")
    print(f"📊 Total slides: {len(prs.slides)}")
    return output_file

if __name__ == "__main__":
    create_concise_presentation()
